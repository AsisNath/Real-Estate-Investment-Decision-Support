"""Generate the 5-minute final presentation deck.

Kept as a script rather than a hand-built file so the deck can be regenerated
whenever the project's numbers change - a slide claiming 93 tests when the suite
has 138 is the kind of error that undermines a demo. The figures below are the
ones verified at build time; update them here, not in PowerPoint, then re-run:

    python scripts/build_presentation.py

Output: docs/NorthStar-Presentation.pptx (16:9, speaker notes with timings).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "docs" / "submission" / "NorthStar-Presentation.pptx"
SHOTS = ROOT / "docs" / "screenshots"

# Same palette as the app's Policy Brief, so the deck and the artifact match.
BRAND = RGBColor(0x1D, 0x5C, 0x4D)
INK = RGBColor(0x1F, 0x27, 0x33)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
HIGH = RGBColor(0xB3, 0x26, 0x1E)
LIGHT = RGBColor(0xEA, 0xF3, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD9, 0xE0, 0xE7)
# Mid tint for large decorative numerals - LIGHT disappears when projected.
TINT = RGBColor(0x9C, 0xC0, 0xB4)

TEAM = "Team 5 · Ashish Nath · Justin Kretschman"
REPO = "github.com/AsisNath/Real-Estate-Investment-Decision-Support"

W, H = Inches(13.333), Inches(7.5)


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def _para(frame, text, size, color=INK, bold=False, space_after=8, first=False, align=None):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.text = text
    para.space_after = Pt(space_after)
    if align is not None:
        para.alignment = align
    for run in para.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return para


def _band(slide, top, height, color):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, W, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _slide(prs, title=None, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    if title:
        _band(slide, 0, Inches(0.09), BRAND)
        frame = _textbox(slide, Inches(0.65), Inches(0.42), W - Inches(1.3), Inches(1.1))
        if kicker:
            _para(frame, kicker.upper(), 12, BRAND, bold=True, space_after=2, first=True)
            _para(frame, title, 30, INK, bold=True, space_after=0)
        else:
            _para(frame, title, 30, INK, bold=True, space_after=0, first=True)
    return slide


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _bullets(frame, items, size=16, first=False):
    """Render bullets into `frame`.

    `first` controls whether the first item reuses paragraphs[0]. It defaults to
    False because a text frame always ships with one empty paragraph, and reusing
    it silently overwrote any heading written into the frame beforehand - which
    is exactly how three column headings vanished from the rendered deck.
    """
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            label, detail = item
            para = _para(frame, f"{label}  ", size, INK, bold=True, space_after=3, first=first and i == 0)
            run = para.add_run()
            run.text = detail
            run.font.size = Pt(size)
            run.font.color.rgb = MUTED
            run.font.name = "Segoe UI"
        else:
            _para(frame, item, size, INK, space_after=9, first=first and i == 0)


def _picture(slide, filename, left, top, width):
    """Drop a screenshot on the slide, scaled to `width`.

    A pitch deck needs product evidence, not a description of the product. These
    are real captures of the running app, so the audience sees the artifact even
    if the live demo is rushed or fails.
    """
    path = SHOTS / filename
    if not path.exists():
        return None
    shape = slide.shapes.add_picture(str(path), left, top, width=width)
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.75)
    return shape


def _stat_card(slide, left, top, width, number, label, accent=BRAND):
    from pptx.enum.shapes import MSO_SHAPE

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.35))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RULE
    card.shadow.inherit = False
    card.text_frame.word_wrap = True
    frame = card.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.12)
    _para(frame, number, 30, accent, bold=True, space_after=0, first=True, align=PP_ALIGN.CENTER)
    _para(frame, label, 11, MUTED, space_after=0, align=PP_ALIGN.CENTER)


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------------------------------------------------------------- 1. Title
    slide = _slide(prs)
    _band(slide, 0, H, BRAND)
    frame = _textbox(slide, Inches(1.0), Inches(2.05), W - Inches(2.0), Inches(3.4))
    _para(frame, "BUKD-X500 · AGENTIC AI SYSTEMS · DESIGN STUDIO", 13, LIGHT, bold=True, space_after=14, first=True)
    _para(frame, "NorthStar", 52, WHITE, bold=True, space_after=0)
    _para(frame, "Property Investment Consulting", 30, LIGHT, space_after=22)
    _para(frame, "“Should I buy this property?”", 22, WHITE, bold=True, space_after=6)
    _para(
        frame,
        "Answered with auditable numbers and current local rules.",
        17,
        LIGHT,
        space_after=26,
    )
    _para(frame, TEAM, 15, WHITE, bold=True, space_after=4)
    _para(frame, REPO, 13, LIGHT, space_after=0)
    _notes(
        slide,
        "[0:00-0:20] Two-sentence hook.\n"
        "NorthStar answers one question for a small real-estate investor: should I buy "
        "this property? The numbers are deterministic Python. The local rules come from "
        "AI research that runs itself and cites every source.",
    )

    # ------------------------------------------------------------- 2. Problem
    slide = _slide(prs, "The same question, two answers, two ZIP codes apart", kicker="The problem")
    frame = _textbox(slide, Inches(0.65), Inches(1.62), Inches(6.35), Inches(4.4))
    _bullets(
        frame,
        [
            "A small investor decides with a spreadsheet plus whatever Google returns.",
            "Rental rules are hyper-local, change often, and the top search results are "
            "property-management blogs that go stale silently.",
            "Model short-term-rental income on a property where STR is illegal and every "
            "number downstream is wrong - IRR, cash flow, the recommendation.",
        ],
        size=16,
        first=True,
    )
    from pptx.enum.shapes import MSO_SHAPE

    # The lead-in belongs directly above the cards it introduces. On the left it
    # dangled with empty space beneath, which read as a missing bullet.
    label = _textbox(slide, Inches(7.25), Inches(1.6), Inches(5.4), Inches(0.4))
    _para(label, "We found this in our own test data:", 15, INK, bold=True, space_after=0, first=True)

    for i, (zipc, city, verdict, color) in enumerate(
        [
            ("63043", "Maryland Heights, MO", "STR allowed for investors\nowner OR manager within 1 hour", BRAND),
            ("63021", "Ballwin, MO", "STR closed to investors\nowner must live there 180 days/yr", HIGH),
        ]
    ):
        left = Inches(7.25)
        top = Inches(2.12) + Inches(1.62) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.4), Inches(1.42))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.shadow.inherit = False
        cf = card.text_frame
        cf.word_wrap = True
        cf.margin_left = Inches(0.2)
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Autoshapes centre their text by default; force left so the heading and
        # the body line up with each other.
        _para(cf, f"{zipc} · {city}", 14, color, bold=True, space_after=3, first=True, align=PP_ALIGN.LEFT)
        _para(cf, verdict, 13, INK, space_after=0, align=PP_ALIGN.LEFT)

    frame = _textbox(slide, Inches(7.25), Inches(5.45), Inches(5.4), Inches(1.5))
    _para(
        frame,
        "Same state. Same county. Same statutes.\nOpposite conclusions.",
        16,
        HIGH,
        bold=True,
        space_after=6,
        first=True,
    )
    _para(
        frame,
        "No static dataset catches that. It has to be researched, per address, and kept fresh.",
        13,
        MUTED,
        space_after=0,
    )
    _notes(
        slide,
        "[0:20-1:10] Lead with the proof, not the premise.\n"
        "Maryland Heights lets an investor run an STR if the owner OR a manager is within "
        "an hour. Ballwin, two ZIPs away, requires the owner to live there 180 days a year "
        "- so a pure investment property cannot do it at all. Same county, same state law, "
        "opposite answers. If you underwrite Ballwin using Maryland Heights' rule, your "
        "entire revenue model is fiction. That is the gap we built for.",
    )

    # ------------------------------------------------------------- 3. Solution
    slide = _slide(prs, "Two layers that never mix", kicker="What we built")
    frame = _textbox(slide, Inches(0.65), Inches(1.7), Inches(5.5), Inches(4.6))
    _para(frame, "1 · The math", 19, BRAND, bold=True, space_after=2, first=True)
    _para(frame, "Plain Python. Unit-tested. No model touches a number.", 15, MUTED, space_after=16)
    _para(frame, "2 · The rules", 19, BRAND, bold=True, space_after=2)
    _para(frame, "An Agent Skill researches local law with live web search\nand cites every source.", 15, MUTED, space_after=16)
    _para(frame, "They meet at the decision", 19, INK, bold=True, space_after=2)
    _para(frame, "Findings become risk flags, diligence items,\nand conflicts with your own assumptions.", 15, MUTED, space_after=14)
    _para(frame, "Research never blocks the report.", 15, HIGH, bold=True, space_after=0)

    _picture(slide, "02-investor-report.png", Inches(6.55), Inches(1.35), Inches(6.15))
    _notes(
        slide,
        "[1:10-1:55] The one-sentence architecture. The screenshot is the point - let them look.\n"
        "Two layers that never mix. The financial model is ordinary Python - we deliberately "
        "kept the LLM out of the arithmetic, because a hallucinated IRR is worse than no IRR. "
        "The AI does the part it is actually good at: reading the web and citing sources. "
        "They meet only at the decision layer.",
    )

    # ------------------------------------------------- 4. Design decisions
    slide = _slide(prs, "Four decisions we would defend", kicker="Design & approach")
    for i, (num, head, sub) in enumerate(
        [
            ("1", "Keep the LLM out of the arithmetic", "A hallucinated IRR is worse than none"),
            ("2", "Provenance enforced in code", "AI-written and human-written notes cannot mix"),
            ("3", "Use the agent the user already signed into", "Zero credential setup, no key stored"),
            ("4", "State missing data, never fill the gap", "“Confirm with the agency”, not a blog"),
        ]
    ):
        top = Inches(1.85) + Inches(1.28) * i
        num_frame = _textbox(slide, Inches(0.65), top - Inches(0.12), Inches(0.7), Inches(0.8))
        _para(num_frame, num, 34, TINT, bold=True, space_after=0, first=True)
        frame = _textbox(slide, Inches(1.35), top, W - Inches(2.0), Inches(1.1))
        _para(frame, head, 20, INK, bold=True, space_after=1, first=True)
        _para(frame, sub, 14, MUTED, space_after=0)
    _notes(
        slide,
        "[1:55-2:25] Only 30 seconds. Say TWO of the four out loud - 1 and 3 land best - and let the slide carry the others.\n"
        "Decision 3 is the one we are proudest of. Our first build made the app call an API "
        "directly, which meant every user needed their own key. We realised the Skill is just "
        "instructions - so we run it inside a CLI the user is already logged into. Setup went "
        "from 'get an API key' to nothing at all.",
    )

    # ------------------------------------------------------------- 5. Demo
    slide = _slide(prs, "Live demo", kicker="Working prototype")
    frame = _textbox(slide, Inches(0.65), Inches(1.6), Inches(6.0), Inches(1.4))
    _para(frame, "Same property, before and after research", 17, INK, bold=True, space_after=8, first=True)

    from pptx.enum.shapes import MSO_SHAPE

    for i, (title, rows, color) in enumerate(
        [
            ("BEFORE", ["STR income: assumed viable", "Policy risk:  medium", "2 high-attention risks"], MUTED),
            ("AFTER", ["STR income: unavailable", "Policy risk:  HIGH", "5 high-attention risks"], HIGH),
        ]
    ):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65) + Inches(3.15) * i, Inches(2.35),
            Inches(2.9), Inches(1.95),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.shadow.inherit = False
        cf = card.text_frame
        cf.word_wrap = True
        cf.margin_left = Inches(0.18)
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _para(cf, title, 13, color, bold=True, space_after=5, first=True, align=PP_ALIGN.LEFT)
        for row in rows:
            _para(cf, row, 13, INK, space_after=3, align=PP_ALIGN.LEFT)

    frame = _textbox(slide, Inches(0.65), Inches(4.55), Inches(6.0), Inches(2.4))
    _para(frame, "Demo path", 15, BRAND, bold=True, space_after=5, first=True)
    for step in [
        "Analyze  →  report is instant and complete",
        "New ZIP  →  research starts itself, unprompted",
        "Saved note  →  every fact sourced, dated, tagged",
        "One click  →  printable Policy Brief for a client",
    ]:
        _para(frame, step, 14, MUTED, space_after=5)

    _picture(slide, "04-policy-brief.png", Inches(7.15), Inches(1.5), Inches(5.55))
    _notes(
        slide,
        "[2:25-4:15] The longest block - 110 seconds. Pre-load the app; do not set up on "
        "stage, and have the recorded backup queued.\n"
        "Say the before/after out loud: before research the model assumed short-term rental "
        "income; after, the app flags it unavailable, policy risk goes medium to HIGH, and "
        "high-attention risks go from two to five.\n"
        "Do NOT say research flipped the recommendation - on this property it was already "
        "Reject on cash flow. The honest claim is that the reasoning changed, and on a "
        "marginal deal it would flip the decision.\n"
        "If time is short, cut the Policy Brief step and keep the research-runs-itself moment.\n"
        "Worth 15 seconds if it fits: a note in our repo for Wildwood 63038 was researched by "
        "the app unattended, 20 official citations, nobody triggered it by hand.",
    )

    # -------------------------------------------------- 6. Testing & security
    slide = _slide(prs, "Tested and hardened, not just working", kicker="Quality & security")
    for i, (number, label) in enumerate(
        [("138", "automated tests"), ("41,488", "US ZIPs bundled\nfor offline checks"),
         ("10", "cited policy notes"), ("0", "secrets in the repo")]
    ):
        _stat_card(slide, Inches(0.65) + Inches(3.13) * i, Inches(1.55), Inches(2.9), number, label)

    frame = _textbox(slide, Inches(0.65), Inches(3.25), Inches(5.9), Inches(3.4))
    _para(frame, "How we tested", 16, BRAND, bold=True, space_after=8, first=True)
    _bullets(
        frame,
        [
            "Financial model unit-tested against hand-computed values.",
            "Two near-opposite jurisdictions used as the research test: Austin and Los Angeles.",
            "Tests force research off, so running the suite can never spend real money.",
            "Bugs the tests caught: a lock deadlock that hung the server, and a stale note "
            "that blocked its own refresh.",
        ],
        size=13.5,
    )

    frame = _textbox(slide, Inches(6.95), Inches(3.25), Inches(5.75), Inches(3.4))
    _para(frame, "Security decisions", 16, BRAND, bold=True, space_after=8, first=True)
    _bullets(
        frame,
        [
            "Path traversal blocked - notes cannot escape the knowledge bank.",
            "Note text is escaped before rendering, so a document cannot inject markup "
            "into a brief a client opens.",
            "The research agent runs read-only; the app decides what reaches disk, so a "
            "reply that is not a policy note is never saved.",
            "Secrets live in a git-ignored .env; the app runs fully offline without one.",
        ],
        size=13.5,
    )
    _notes(
        slide,
        "[4:15-4:45] Do not read the cards; point at them.\n"
        "Threat model: this is a local app, so the realistic risk is not a remote attacker - "
        "it is untrusted text. Policy notes come from the open web and from PDFs a user drops "
        "in, and those get rendered into a document handed to a client. So we escape "
        "everything and validate before saving.",
    )

    # -------------------------------------------------------- 7. Reflections
    slide = _slide(prs, "What we learned", kicker="Reflections")
    frame = _textbox(slide, Inches(0.65), Inches(1.62), Inches(3.85), Inches(5.0))
    _para(frame, "Worked", 16, BRAND, bold=True, space_after=8, first=True)
    _bullets(
        frame,
        [
            "Writing standards instead of steps. The Skill says what good research looks "
            "like; the agent works out how.",
            "Making provenance a code boundary rather than a naming convention.",
            "Keeping the AI out of the numbers.",
        ],
        size=13.5,
    )

    frame = _textbox(slide, Inches(4.75), Inches(1.62), Inches(3.85), Inches(5.0))
    _para(frame, "Did not", 16, HIGH, bold=True, space_after=8, first=True)
    _bullets(
        frame,
        [
            "First research design needed an API key per user. Wrong shape - the Skill "
            "already ran inside tools people had.",
            "A stale server made fixed code look broken for an hour.",
            "Documents users supply are often scans; the parser silently ignored them "
            "until we surfaced it.",
        ],
        size=13.5,
    )

    frame = _textbox(slide, Inches(8.85), Inches(1.62), Inches(3.85), Inches(5.0))
    _para(frame, "Next", 16, MUTED, bold=True, space_after=8, first=True)
    _bullets(
        frame,
        [
            "OCR so scanned HOA declarations parse themselves.",
            "Comparison view across several candidate properties.",
            "Scheduled re-research so notes refresh before they age out.",
        ],
        size=13.5,
    )

    _band(slide, H - Inches(0.95), Inches(0.95), BRAND)
    frame = _textbox(slide, Inches(0.65), H - Inches(0.78), W - Inches(1.3), Inches(0.6))
    _para(frame, f"{REPO}   ·   {TEAM}", 13, WHITE, bold=True, space_after=0, first=True, align=PP_ALIGN.CENTER)
    _notes(
        slide,
        "[4:45-5:00] Close on the lesson, not a summary.\n"
        "The biggest shift: we spent our time writing standards - what counts as a verified "
        "fact, when to say 'unverified' - rather than wiring steps together. The agent handled "
        "the how, including recovering when a government site blocked it. That is a different "
        "job from programming a pipeline.",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path.relative_to(path.parent.parent)}  ({path.stat().st_size / 1024:.0f} KB)")
